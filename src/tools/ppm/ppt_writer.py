from copy import deepcopy
from datetime import datetime, date
import re
from typing import List, Tuple
import unicodedata

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
import textwrap



from openpyxl import load_workbook
COMPACT_MIN_ROW_H = 12000           # ~0.13" in EMU
LABEL_INSIDE_THRESHOLD = 220000     # bar width threshold to place label inside

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_EMBED = f"{{{_REL_NS}}}embed"
_LINK  = f"{{{_REL_NS}}}link"

def _norm_etapa_key(etapa_key: str) -> str:
    if etapa_key is None:
        return ""
    s = str(etapa_key).strip().lower()
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))  # quita acentos
    s = s.replace(" ", "_")
    s = re.sub(r"[^a-z0-9_]+", "_", s)
    s = re.sub(r"_+", "_", s).strip("_")
    return s

def _to_percent_number(x) -> float:
    """
    Normaliza avance_* a un número en escala 0..100
    - Si viene 0.96 -> 96.0
    - Si viene 96 -> 96.0
    - Si viene "96%" -> 96.0
    """
    if x is None or x == "":
        return 0.0
    if isinstance(x, str):
        s = x.strip().replace("%", "").replace(",", ".")
        try:
            return float(s)
        except Exception:
            return 0.0
    try:
        v = float(x)
    except Exception:
        return 0.0
    # Si está en 0..1, es fracción
    if 0.0 <= v <= 1.0:
        return v * 100.0
    return v


def _color_estado_general_from_avances(sem_shape, avance_planeado, avance_real):
    """
    diff = planeado - real (en puntos porcentuales)
    <= 5  -> verde
    6..10 -> amarillo
    > 10  -> rojo
    """
    if not sem_shape:
        return

    p = _to_percent_number(avance_planeado)
    r = _to_percent_number(avance_real)
    diff = p - r

    sem_shape.fill.solid()
    if diff <= 5:
        sem_shape.fill.fore_color.rgb = RGBColor(0, 176, 80)      # verde
    elif diff <= 10:
        sem_shape.fill.fore_color.rgb = RGBColor(255, 192, 0)     # amarillo
    else:
        sem_shape.fill.fore_color.rgb = RGBColor(192, 0, 0)       # rojo





def _estimate_lines(text: str, max_chars: int) -> int:
    """Heurística de wrap por caracteres."""
    if not text:
        return 1
    lines = textwrap.wrap(text, width=max(10, max_chars), break_long_words=False, replace_whitespace=False)
    return max(1, len(lines))



def _set_descripcion_estatus_with_checks(slide, shp_text, raw_text: str, marker="(*)", src_icon_name=None):
    """
    Inserta un check INLINE (✓) al inicio de cada "párrafo" (línea) y,
    SI hay contenido, agrega arriba el título:

        Desarrollo SYP:   (bold)
        <línea en blanco>
        ✓ ...contenido...

    Reglas:
    - Si el texto contiene el marcador en AL MENOS una línea, SOLO esas líneas llevan check.
    - Si no hay marcador en ninguna línea, todas las líneas no vacías llevan check.
    - El marcador se elimina del texto final, tolerando variantes: '(*)', '(* )', '( * )', etc.
    """
    if not shp_text or not getattr(shp_text, "has_text_frame", False):
        return

    tf = shp_text.text_frame

    # Mantén settings del textbox lo más posible
    _disable_textbox_autofit_xml(shp_text)
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.word_wrap = True

    # Captura estilo "donor" del primer run existente (si hay)
    donor_font = None
    donor_align = None
    donor_color = None
    try:
        donor_vanchor = tf.vertical_anchor
    except Exception:
        donor_vanchor = None

    donor_margins = (
        getattr(tf, "margin_left", None),
        getattr(tf, "margin_right", None),
        getattr(tf, "margin_top", None),
        getattr(tf, "margin_bottom", None),
    )
    try:
        p0 = tf.paragraphs[0]
        donor_align = p0.alignment
        if p0.runs:
            donor_font = p0.runs[0].font
            try:
                donor_color = p0.runs[0].font.color.rgb
            except Exception:
                donor_color = None
    except Exception:
        pass

    text = "" if raw_text is None else str(raw_text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Tratamos cada LÍNEA no vacía como párrafo (es lo que estás usando en Excel)
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]

    # Si no hay contenido, no pintes nada (ni título)
    if not lines:
        try:
            tf.clear()
        except Exception:
            pass
        return

    # Marcador tolerante: ( * ) con espacios opcionales
    marker_re = re.compile(r"\(\s*\*\s*\)\s*")

    has_any_marker = any(marker_re.search(ln) is not None for ln in lines)

    paragraphs = []
    marked = []
    for ln in lines:
        has = marker_re.search(ln) is not None
        cleaned = marker_re.sub("", ln).strip()
        paragraphs.append(cleaned)
        marked.append(has)

    if not has_any_marker:
        marked = [True for _ in paragraphs]

    # Rebuild text frame (control total de runs/párrafos)
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    # Restituye anchor/margins si estaban
    try:
        if donor_vanchor is not None:
            tf.vertical_anchor = donor_vanchor
    except Exception:
        pass
    try:
        ml, mr, mt, mb = donor_margins
        if ml is not None:
            tf.margin_left = ml
        if mr is not None:
            tf.margin_right = mr
        if mt is not None:
            tf.margin_top = mt
        if mb is not None:
            tf.margin_bottom = mb
    except Exception:
        pass

    master_pt = _get_effective_master_size_pt(shp_text)

    # -------------------------------
    # TÍTULO + LÍNEA EN BLANCO
    # -------------------------------
    p_title = tf.paragraphs[0]
    if donor_align is not None:
        p_title.alignment = donor_align

    r_title = p_title.add_run()
    r_title.text = "Desarrollo SYP:"
    if donor_font is not None:
        _clone_font_no_color(r_title.font, donor_font)
    if master_pt is not None:
        r_title.font.size = Pt(master_pt)
    r_title.font.bold = True
    try:
        r_title.font.color.rgb = RGBColor(0, 0, 0)
    except Exception:
        pass

    # Doble salto de línea (1 párrafo vacío)
    p_space = tf.add_paragraph()
    if donor_align is not None:
        p_space.alignment = donor_align
    p_space.text = ""

    # -------------------------------
    # CONTENIDO (checks + texto)
    # -------------------------------
    for ptext, put_check in zip(paragraphs, marked):
        p = tf.add_paragraph()
        if donor_align is not None:
            p.alignment = donor_align

        if put_check:
            r_chk = p.add_run()
            r_chk.text = "✓  "
            if donor_font is not None:
                _clone_font(r_chk.font, donor_font)
            if master_pt is not None:
                r_chk.font.size = Pt(master_pt)
            r_chk.font.color.rgb = RGBColor(0, 176, 80)

        r_txt = p.add_run()
        r_txt.text = ptext
        if donor_font is not None:
            # No copies color: en tu plantilla el texto puede venir blanco por tema/estilo
            _clone_font_no_color(r_txt.font, donor_font)
        if master_pt is not None:
            r_txt.font.size = Pt(master_pt)
        # Texto SIEMPRE negro (y el check verde)
        try:
            r_txt.font.color.rgb = RGBColor(0, 0, 0)
        except Exception:
            pass


def _disable_textbox_autofit_xml(shape):
    """
    Fuerza <a:noAutofit/> en bodyPr para evitar shrink/autoajuste que cambia la fuente.
    """
    if not getattr(shape, "has_text_frame", False):
        return

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    txb = shape._element.find(".//p:txBody", namespaces=ns)
    if txb is None:
        return

    bodyPr = txb.find("./a:bodyPr", namespaces=ns)
    if bodyPr is None:
        return

    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        el = bodyPr.find(f"./{tag}", namespaces=ns)
        if el is not None:
            bodyPr.remove(el)

    from pptx.oxml.xmlchemy import OxmlElement
    bodyPr.append(OxmlElement("a:noAutofit"))

def _get_effective_master_size_pt(shape):
    """
    Devuelve el tamaño "más probable" en pt:
    1) rPr/@sz del PRIMER run (si existe) -> es el override REAL que tú ves como 18
    2) defRPr/@sz de lstStyle (fallback típico del master/layout)
    3) None si no se puede
    """
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    # 1) Primer run rPr sz (override explícito)
    rPr = shape._element.find(".//a:p/a:r/a:rPr", namespaces=ns)
    if rPr is not None:
        sz = rPr.get("sz")
        if sz and str(sz).isdigit():
            return int(sz) / 100.0

    # 2) lstStyle defRPr sz (estilo del placeholder)
    defrpr = shape._element.find(".//a:lstStyle/a:lvl1pPr/a:defRPr", namespaces=ns)
    if defrpr is not None:
        sz = defrpr.get("sz")
        if sz and str(sz).isdigit():
            return int(sz) / 100.0

    # 3) pPr defRPr (otra variante)
    defrpr = shape._element.find(".//a:p/a:pPr/a:defRPr", namespaces=ns)
    if defrpr is not None:
        sz = defrpr.get("sz")
        if sz and str(sz).isdigit():
            return int(sz) / 100.0

    return None

def _norm_status(x) -> str:
    # Normaliza estatus (Excel suele traer None, espacios, etc.)
    if x is None:
        return ""
    s = str(x).strip().upper()
    if s in ("", "NONE", "NAN", "N/A", "NA", "-"):
        return ""
    s = s.replace(" ", "_")
    return s


def _apply_stage_icon(slide, etapa_key: str, status: str):
    """
    status:
      COMPLETO -> check visible
      EN_CURSO -> gear visible
      vacío / otro -> no mostrar icono

    Nota:
    - No asumimos que existan ambos shapes.
    - Siempre ocultamos primero para evitar iconos "pegados".
    """

    etapa_key = _norm_etapa_key(etapa_key)  # 🔥 CLAVE
    st = _norm_status(status)

    shp_check = _find_shape(slide, f"ico_etapa_{etapa_key}_check")
    shp_gear  = _find_shape(slide, f"ico_etapa_{etapa_key}_gear")

    # Si no existe ninguno, no hay nada que hacer
    if not shp_check and not shp_gear:
        return

    # Reset duro: ocultar todo lo que exista
    if shp_check:
        _set_shape_hidden(shp_check, True)
    if shp_gear:
        _set_shape_hidden(shp_gear, True)

    # Vacío => no mostrar nada
    if not st:
        return

    if st == "COMPLETO":
        if shp_check:
            _set_shape_hidden(shp_check, False)
        return

    if st == "EN_CURSO":
        if shp_gear:
            _set_shape_hidden(shp_gear, False)
        return

    # Cualquier otro valor => ninguno (ya ocultos)


def _force_visible(shape):
    # 1) API (si existe)
    try:
        shape.hidden = False
    except Exception:
        pass

    # 2) XML (por si quedó marcado como hidden en cNvPr)
    try:
        cNvPr = shape._element.xpath(".//*[local-name()='cNvPr']")[0]
        if cNvPr.get("hidden") is not None:
            cNvPr.attrib.pop("hidden", None)
    except Exception:
        pass


def _set_shape_hidden(shape, hide: bool):
    """Oculta/muestra un shape vía XML (más confiable que .visible/.hidden con master/layout)."""
    try:
        cNvPr = shape._element.xpath(".//*[local-name()='cNvPr']")[0]
    except Exception:
        return
    if hide:
        cNvPr.set("hidden", "1")
    else:
        cNvPr.attrib.pop("hidden", None)



def _set_etapas_icons(slide, gantt_rows):
    """
    Lee gantt_rows (hoja GANTT) y actualiza iconos.
    Se toma la etapa por actividad_id.
    """
    # mapeo actividad_id (excel) -> sufijo de nombre en ppt
    # AJUSTA si tus nombres difieren
    map_etapa = {
        "ESTIMACION": "estimacion",
        "PLANEACION": "planeacion",
        "ANALISIS": "analisis_tecnico",
        "DISENO": "diseno_detallado",
        "REALIZACION": "realizacion",
        "QA": "qa",
        "IMPL": "implementacion",
        "GARANTIA": "garantia",
    }

    # default: completo si no viene campo
    for g in gantt_rows:
        # Prefer actividad_id if provided; fall back to actividad text
        act_id = (_gget(g, 'actividad_id', None) or '').strip().upper()
        if not act_id:
            act_txt = (_gget(g, 'actividad', '') or '').strip()
            act_id = _strip_accents(act_txt).strip().upper().replace(' ', '_')
        etapa_key = map_etapa.get(act_id)
        if not etapa_key:
            continue

        status = _gget(g, "estatus_etapa", None)  # <- columna en Excel
        if status is None or str(status).strip() == "":
            status = ""  # vacío => no icono

        _apply_stage_icon(slide, etapa_key, status)


def _gget(g, key, default=None):
    if g is None:
        return default
    if isinstance(g, dict):
        return g.get(key, default)
    return getattr(g, key, default)

def _get_master_font_size_pt_from_textbox(shape):
    """
    Intenta obtener el tamaño de fuente "real" del master/layout para este textbox:
    busca defRPr/@sz en lstStyle (1/100 pt).
    Devuelve float (pt) o None si no existe.
    """
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

    txb = shape._element.find(".//p:txBody", namespaces=ns)
    if txb is None:
        return None

    # 1) Preferencia: lstStyle/lvl1pPr/defRPr (lo más típico en cuadros de texto del master)
    defrpr = txb.find(".//a:lstStyle/a:lvl1pPr/a:defRPr", namespaces=ns)
    if defrpr is not None:
        sz = defrpr.get("sz")
        if sz and str(sz).isdigit():
            return int(sz) / 100.0

    # 2) Alternativa: primer párrafo tiene pPr/defRPr
    defrpr = txb.find(".//a:p/a:pPr/a:defRPr", namespaces=ns)
    if defrpr is not None:
        sz = defrpr.get("sz")
        if sz and str(sz).isdigit():
            return int(sz) / 100.0

    return None

def _to_date(x):
    if x is None or x == "":
        return None
    if isinstance(x, datetime):
        return x.date()
    if isinstance(x, date):
        return x
    try:
        if hasattr(x, "to_pydatetime"):
            return x.to_pydatetime().date()
    except Exception:
        pass
    try:
        return datetime.fromisoformat(str(x)).date()
    except Exception:
        return None


def _clone_font(dst_font, src_font):
    if src_font is None:
        return
    if src_font.size is not None:
        dst_font.size = src_font.size
    if src_font.name is not None:
        dst_font.name = src_font.name
    try:
        if src_font.color is not None and src_font.color.rgb is not None:
            dst_font.color.rgb = src_font.color.rgb
    except Exception:
        pass


def _clone_font_no_color(dst_font, src_font):
    """Copia SOLO size/name (y opcionalmente bold/italic/underline), pero nunca copia color."""
    if src_font is None:
        return
    if src_font.size is not None:
        dst_font.size = src_font.size
    if src_font.name is not None:
        dst_font.name = src_font.name
    for attr in ("bold", "italic", "underline"):
        v = getattr(src_font, attr, None)
        if v is not None:
            setattr(dst_font, attr, v)





def _clone_font_full(dst_font, src_font):
    if src_font is None:
        return
    for attr in ("size", "name", "bold", "italic", "underline"):
        v = getattr(src_font, attr, None)
        if v is not None:
            setattr(dst_font, attr, v)
    try:
        if src_font.color is not None and src_font.color.rgb is not None:
            dst_font.color.rgb = src_font.color.rgb
    except Exception:
        pass


def _append_table_row_clone(table, copy_from_row_idx=0):
    tbl = table._tbl
    tr_lst = getattr(tbl, "tr_lst", None)
    if not tr_lst:
        return
    src_idx = max(0, min(copy_from_row_idx, len(tr_lst) - 1))
    src_tr = tr_lst[src_idx]
    new_tr = deepcopy(src_tr)
    for tc in getattr(new_tr, "tc_lst", []):
        for t_el in tc.xpath(".//*[local-name()='t']"):
            t_el.text = ""
    tbl.append(new_tr)


def _ensure_table_rows_real(table, total_rows_needed, copy_from_row_idx=0):
    while len(table.rows) < total_rows_needed:
        _append_table_row_clone(table, copy_from_row_idx=copy_from_row_idx)


def _get_table_donor_run(table, prefer_row=1, prefer_col=0):
    rows = len(table.rows)
    cols = len(table.columns)
    for r, c in [(prefer_row, prefer_col), (0, prefer_col), (0, 0)]:
        if 0 <= r < rows and 0 <= c < cols:
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                if p.runs:
                    return p.runs[0]
    for r in range(rows):
        for c in range(cols):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                if p.runs:
                    return p.runs[0]
    return None

def _set_text_keep_style(shape, text: str):
    """
    Reemplaza texto en un TEXTBOX/PLACEHOLDER respetando el master:
    - No borra XML (no remove de runs/párrafos)
    - Desactiva AutoFit real (XML) para que PPT no cambie tamaño
    - "Ancla" el size al que el master/override realmente tenía (evita salto 18 -> 32)
    - Permite saltos de línea sin usar \v (evita _x000B_)
    """
    if not getattr(shape, "has_text_frame", False):
        return

    tf = shape.text_frame
    text = "" if text is None else str(text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # 1) No-autofit
    _disable_textbox_autofit_xml(shape)
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.word_wrap = True

    # 2) Obtener tamaño efectivo del master/override y fijarlo explícitamente
    master_pt = _get_effective_master_size_pt(shape)

    # 3) Asegurar párrafo + run base EXISTENTE
    p0 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    r0 = p0.runs[0] if p0.runs else p0.add_run()

    # Fijar size ANTES de cambiar texto (evita que PPT re-aplique 32 del placeholder)
    if master_pt is not None:
        r0.font.size = Pt(master_pt)

    # 4) Escribir TODO en el primer run (no crear párrafos nuevos)
    r0.text = text

    # 5) Vaciar cualquier resto sin borrar estructura
    for r in list(p0.runs)[1:]:
        r.text = ""
    for p in list(tf.paragraphs)[1:]:
        for r in p.runs:
            r.text = ""

def _set_cell_text_keep_style(cell, text, donor_run=None, donor_full=False):
    tf = cell.text_frame
    text = "" if text is None else str(text)
    if tf.paragraphs and tf.paragraphs[0].runs:
        p0 = tf.paragraphs[0]
        p0.runs[0].text = text
        for r in list(p0.runs)[1:]:
            p0._p.remove(r._r)
        for p in list(tf.paragraphs)[1:]:
            tf._txBody.remove(p._p)
        return
    p0 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    run = p0.add_run()
    run.text = text
    if donor_run is not None:
        if donor_full:
            _clone_font_full(run.font, donor_run.font)
        else:
            _clone_font(run.font, donor_run.font)


def _fmt_date(d) -> str:
    if d is None:
        return ""
    if isinstance(d, datetime):
        d = d.date()
    if isinstance(d, date):
        return d.strftime("%d-%b")
    return str(d)


def _fmt_money(x) -> str:
    if x is None or x == "":
        return ""
    try:
        return "${:,.2f}".format(float(x))
    except Exception:
        return str(x)


def _find_shape(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None

def _norm_status(x) -> str:
    # Normaliza estatus (Excel suele traer None, espacios, etc.)
    if x is None:
        return ""
    s = str(x).strip().upper()
    if s in ("", "NONE", "NAN", "N/A", "NA", "-"):
        return ""
    s = s.replace(" ", "_")
    return s


def _apply_stage_icon(slide, etapa_key: str, status: str):
    """
    status:
      COMPLETO -> check visible
      EN_CURSO -> gear visible
      vacío / otro -> no mostrar icono
    Requiere que en la plantilla existan:
      ico_etapa_<etapa_key>_check
      ico_etapa_<etapa_key>_gear
    """
    st = _norm_status(status)

    shp_check = _find_shape(slide, f"ico_etapa_{etapa_key}_check")
    shp_gear  = _find_shape(slide, f"ico_etapa_{etapa_key}_gear")

    if not shp_check or not shp_gear:
        return

    # default: ocultar ambos
    _set_shape_hidden(shp_check, True)
    _set_shape_hidden(shp_gear, True)

    if not st:
        return

    if st == "COMPLETO":
        _set_shape_hidden(shp_check, False)
    elif st == "EN_CURSO":
        _set_shape_hidden(shp_gear, False)
    else:
        return


    shp_check = _find_shape(slide, f"ico_etapa_{etapa_key}_check")
    shp_gear  = _find_shape(slide, f"ico_etapa_{etapa_key}_gear")

    # Si tu plantilla aún no tiene ambos, aquí no hay magia: hay que crearlos/duplicarlos en el PPT.
    if not shp_check or not shp_gear:
        return

    if st == "COMPLETO":
        shp_check.visible = True
        shp_gear.visible = False
    else:  # EN_CURSO (o cualquier otro)
        shp_check.visible = False
        shp_gear.visible = True


def _set_etapas_icons(slide, gantt_rows):
    """
    Lee gantt_rows (hoja GANTT) y actualiza iconos.
    Se toma la etapa por actividad_id.
    """
    # mapeo actividad_id (excel) -> sufijo de nombre en ppt
    # AJUSTA si tus nombres difieren
    map_etapa = {
        "ESTIMACION": "estimacion",
        "PLANEACION": "planeacion",
        "ANALISIS": "analisis_tecnico",
        "DISENO": "diseno_detallado",
        "REALIZACION": "realizacion",
        "QA": "qa",
        "IMPL": "implementacion",
        "GARANTIA": "garantia",
    }

    # default: completo si no viene campo
    for g in gantt_rows:
        # Prefer actividad_id if provided; fall back to actividad text
        act_id = (_gget(g, 'actividad_id', None) or '').strip().upper()
        if not act_id:
            act_txt = (_gget(g, 'actividad', '') or '').strip()
            act_id = _strip_accents(act_txt).strip().upper().replace(' ', '_')
        etapa_key = map_etapa.get(act_id)
        if not etapa_key:
            continue

        status = _gget(g, "estatus_etapa", None)  # <- columna en Excel
        if status is None or str(status).strip() == "":
            status = ""  # vacío => no icono

        _apply_stage_icon(slide, etapa_key, status)


def _delete_generated(slide, prefix: str):
    to_delete = [shp for shp in slide.shapes if shp.name.startswith(prefix)]
    for shp in to_delete:
        el = shp._element
        el.getparent().remove(el)



def _clone_shape_to_slide(slide, src_shape, new_name: str, gen_prefix="gen_gantt_"):
    """
    Clona un shape (oculto o no) desde layout/master/slide al slide actual.
    - Copia XML
    - Remapea rIds si el shape trae recursos (imágenes, etc.)
    - Fuerza visibilidad del clon
    """
    icon_el = deepcopy(src_shape._element)
    try:
        _remap_all_rids_to_slide_part(slide, icon_el)
    except Exception:
        pass

    slide.shapes._spTree.insert_element_before(icon_el, "p:extLst")
    cloned = next((sh for sh in slide.shapes if sh._element is icon_el), None)
    if not cloned:
        return None

    # Nombrado consistente
    cloned.name = f"{gen_prefix}{new_name}" if not str(new_name).startswith(gen_prefix) else str(new_name)

    # Fuerza visible aunque el master esté oculto
    _force_visible(cloned)
    try:
        cloned.hidden = False
    except Exception:
        pass

    _bring_to_front(slide, cloned)
    return cloned


def _add_globo_fecha(
    slide,
    x_end: int,
    bar_top: int,
    bar_h: int,
    fecha_txt: str,
    gen_prefix="gen_gantt_",
    gap_px: int = 30,          # <-- más a la derecha "coco"
    scale_to_bar: float = 0.88,
    stretch_x: float = 1.22,   # <-- más largo (solo ancho)
):

    """
    Inserta un globo_fecha clonado y lo posiciona al final de la barra.
    Reglas:
      - El tamaño del globo se deriva del alto de la barra (bar_h) para que "coincida"
      - Se conserva la proporción del globo original (plantilla)
      - Texto se escribe respetando el formato del globo en el master (run existente)
    """
    src = _find_shape_any_master_layout(slide, "globo_fecha")
    if not src:
        return None

    globo = _clone_shape_to_slide(slide, src, f"globo_fecha_{re.sub(r'[^0-9A-Za-z]+','_', str(fecha_txt))}", gen_prefix=gen_prefix)
    if not globo:
        return None

    # 1) Tamaño derivado de la barra: el globo se ajusta proporcionalmente con base en bar_h
    #    (un poco más pequeño para que no se vea "grandote")
    try:
        orig_w, orig_h = int(globo.width), int(globo.height)
        target_h = max(1, int(bar_h * float(scale_to_bar)))
        if orig_h > 0:
            s = target_h / float(orig_h)
            globo.height = int(orig_h * s)
            globo.width  = int(orig_w * s)

            # Estira SOLO el ancho para que se vea más largo "coco"
            globo.width = int(globo.width * float(stretch_x))
    except Exception:
        pass

    # 2) Posición: a la derecha del final de barra, centrado verticalmente con la barra
    try:
        emu_per_px = 9525  # ~96dpi. 1px ≈ 9525 EMU
        globo.left = int(x_end + (gap_px * emu_per_px))
        globo.top  = int(bar_top + (bar_h - globo.height) / 2)
    except Exception:
        pass

    # 3) Texto con formato del master
    if getattr(globo, "has_text_frame", False):
        _set_text_keep_style(globo, fecha_txt)

    globo.name = f"{gen_prefix}globo_{fecha_txt}_{globo.left}_{globo.top}"
    _bring_to_front(slide, globo)
    return globo


def duplicate_slide_full(prs, src_slide):
    dst_slide = prs.slides.add_slide(src_slide.slide_layout)

    # limpia shapes del slide nuevo
    for shp in list(dst_slide.shapes):
        el = shp._element
        el.getparent().remove(el)

    # copia XML de shapes del slide origen
    for shp in src_slide.shapes:
        new_el = deepcopy(shp._element)
        dst_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    src_part = src_slide.part
    dst_part = dst_slide.part

    rid_map = {}

    # remapea TODOS los blips (incluye los que estén dentro de grupos)
    blips = dst_slide._element.xpath(".//a:blip")
    for blip in blips:
        for attr in (
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link",
        ):
            old_rid = blip.get(attr)
            if not old_rid:
                continue

            if old_rid not in rid_map:
                rel = src_part.rels.get(old_rid)
                if rel is None:
                    continue
                new_rid = dst_part.relate_to(rel._target, rel.reltype)
                rid_map[old_rid] = new_rid

            blip.set(attr, rid_map[old_rid])

    return dst_slide


def _month_range(start_dt: date, end_dt: date) -> List[Tuple[int, int]]:
    months = []
    y, m = start_dt.year, start_dt.month
    while (y, m) <= (end_dt.year, end_dt.month):
        months.append((y, m))
        m += 1
        if m == 13:
            m = 1
            y += 1
    return months


def _days_in_month(y: int, m: int) -> int:
    if m == 12:
        nxt = date(y + 1, 1, 1)
    else:
        nxt = date(y, m + 1, 1)
    cur = date(y, m, 1)
    return (nxt - cur).days


def _read_left_styles(left_tbl):
    hdr_cell = left_tbl.cell(0, 0)
    body_cell = left_tbl.cell(1, 0) if len(left_tbl.rows) > 1 else hdr_cell

    def _cell_rgb(cell, fallback=RGBColor(89, 89, 89)):
        try:
            if cell.fill and cell.fill.fore_color and cell.fill.fore_color.rgb:
                return cell.fill.fore_color.rgb
        except Exception:
            pass
        return fallback

    styles = {
        "hdr_rgb": _cell_rgb(hdr_cell, RGBColor(89, 89, 89)),
        "body_rgb": _cell_rgb(body_cell, RGBColor(255, 255, 255)),
        "hdr_run": None,
        "body_run": None,
        "hdr_align": None,
        "body_align": None,
        "hdr_vanchor": None,
        "body_vanchor": None,
    }
    if hdr_cell.text_frame.paragraphs:
        p = hdr_cell.text_frame.paragraphs[0]
        styles["hdr_align"] = p.alignment
        if p.runs:
            styles["hdr_run"] = p.runs[0]
    if body_cell.text_frame.paragraphs:
        p = body_cell.text_frame.paragraphs[0]
        styles["body_align"] = p.alignment
        if p.runs:
            styles["body_run"] = p.runs[0]
    styles["hdr_vanchor"] = getattr(hdr_cell.text_frame, "vertical_anchor", None)
    styles["body_vanchor"] = getattr(body_cell.text_frame, "vertical_anchor", None)
    return styles


def _build_time_grid_in_area(slide, area_shape_name, months, row_heights, left_tbl, gen_prefix="gen_gantt_"):
    area = _find_shape(slide, area_shape_name)
    if not area:
        return None
    L, T, W, H = area.left, area.top, area.width, area.height
    n_months = max(1, len(months))

    left_header_h = left_tbl.rows[0].height
    header_h = left_header_h

    st = _read_left_styles(left_tbl)

    # Header background
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, T, W, header_h)
    hdr.name = f"{gen_prefix}hdr_bg"

    hdr.fill.solid()
    hdr.fill.fore_color.rgb = st["hdr_rgb"]

    hdr.line.fill.background()
    hdr.shadow.inherit = False  # ⬅️ elimina sombra del header

    # Month labels
    col_w = int(W / n_months)
    last_w = W - col_w * (n_months - 1)
    for i, (y, m) in enumerate(months):
        x = L + i * col_w
        w = last_w if i == n_months - 1 else col_w
        tb = slide.shapes.add_textbox(x, T, w, header_h)
        tb.name = f"{gen_prefix}hdr_{y}_{m}"
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = date(y, m, 1).strftime("%b").capitalize()
        if st["hdr_run"]:
            _clone_font_full(r.font, st["hdr_run"].font)
        r.font.color.rgb = RGBColor(255, 255, 255)

        p.alignment = st["hdr_align"] if st["hdr_align"] is not None else PP_ALIGN.CENTER
        try:
            tf.vertical_anchor = st["hdr_vanchor"] if st["hdr_vanchor"] is not None else MSO_VERTICAL_ANCHOR.MIDDLE
        except Exception:
            pass
    # Header vertical separators (white, like tblGantt)
    for i in range(1, n_months):
        x = L + i * col_w
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            T,
            1,
            header_h
        )
        sep.name = f"{gen_prefix}hdr_sep_{i}"

        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(255, 255, 255)  # blanco

        sep.line.fill.background()
        sep.shadow.inherit = False

    # Vertical grid lines
    body_h = sum(row_heights)
    for i in range(1, n_months):
        x = L + i * col_w
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, T + header_h, 1, body_h)
        ln.name = f"{gen_prefix}vline_{i}"

        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(210, 210, 210)  # gris claro

        ln.line.fill.background()
        ln.shadow.inherit = False  # ⬅️ elimina sombra

    # Row guides
    row_tops = []
    y = T + header_h
    for idx, rh in enumerate(row_heights):
        row_tops.append(y)
        gline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, L, int(y + rh - 1), W, 1
        )
        gline.name = f"{gen_prefix}row_{idx}_guide"

        gline.fill.solid()
        gline.fill.fore_color.rgb = RGBColor(225, 225, 225)

        gline.line.fill.background()
        gline.shadow.inherit = False  # ⬅️ elimina sombra

        y += rh

    return {
        "L": L,
        "T": T,
        "header_h": header_h,
        "col_w": col_w,
        "row_tops": row_tops,
        "row_heights": row_heights,
        "months": months,
        "W": W,
        "H": H,
    }


from pptx.util import Pt  # ya lo tienes

def _date_to_x_grid(d: date, grid):
    months = grid["months"]
    L = grid["L"]
    col_w = grid["col_w"]
    if not months:
        return float(L)

    for i, (y, m) in enumerate(months):
        if d.year == y and d.month == m:
            dim = _days_in_month(y, m)

            # Centro del día
            day_pos = max(0.0, min(dim, (d.day - 0.5)))

            fudge_x = 0#int(Pt(18))   # prueba: 1.0, 1.5, 2.0, 3.0
            return L + i * col_w + (day_pos / dim) * col_w + fudge_x

    if (d.year, d.month) < months[0]:
        return float(L)
    return float(L + len(months) * col_w)




def _draw_today_line(slide, grid, gen_prefix="gen_gantt_"):
    today = datetime.today().date()
    months = grid["months"]
    if not months:
        return
    if not ((today.year, today.month) >= months[0] and (today.year, today.month) <= months[-1]):
        return
    x = _date_to_x_grid(today, grid)
    y = grid["T"] + grid["header_h"]
    h = sum(grid["row_heights"])  # body height
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), 2, int(h))
    ln.name = f"{gen_prefix}today"
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor(192, 0, 0)
    ln.line.fill.background()


def _fill_tbl_gantt_left(slide, gantt_rows, gen_prefix="gen_gantt_"):
    shp = _find_shape(slide, "tblGantt")
    if not shp or not shp.has_table:
        return
    t = shp.table
    if len(t.rows) < 2 or len(t.columns) < 2:
        return

    needed_total_rows = 1 + max(1, len(gantt_rows))
    _ensure_table_rows_real(t, needed_total_rows, copy_from_row_idx=1)

    donor = _get_table_donor_run(t, prefer_row=1, prefer_col=0)
    for r in range(1, len(t.rows)):
        _set_cell_text_keep_style(t.cell(r, 0), "", donor)
        _set_cell_text_keep_style(t.cell(r, 1), "", donor)

    max_rows = len(t.rows) - 1
    for i, g in enumerate(gantt_rows[:max_rows]):
        r = i + 1
        _set_cell_text_keep_style(t.cell(r, 0), _gget(g, "actividad", "") or "", donor)
        _set_cell_text_keep_style(t.cell(r, 1), _gget(g, "responsable", "") or "", donor)



def _draw_gantt(slide, gantt_rows, gen_prefix="gen_gantt_"):
    _delete_generated(slide, gen_prefix)

    starts = [_to_date(_gget(g, "fecha_inicio")) for g in gantt_rows]
    ends   = [_to_date(_gget(g, "fecha_fin"))   for g in gantt_rows]
    starts = [d for d in starts if d]
    ends   = [d for d in ends if d]
    if not starts or not ends:
        return

    min_start = min(starts)
    max_end = max(ends)
    months = _month_range(min_start, max_end)
    if not months:
        return

    left_shp = _find_shape(slide, "tblGantt")
    if not left_shp or not left_shp.has_table:
        return
    left_tbl = left_shp.table

    # COMPACT: compute row heights to FIT the available area height (not sum of left rows)
    area = _find_shape(slide, "tblGanttTime")
    if not area:
        return
    left_header_h = left_tbl.rows[0].height
    body_h = max(1, int(area.height - left_header_h))

    n_rows_data = len(gantt_rows)
    each_h = max(COMPACT_MIN_ROW_H, int(body_h / max(1, n_rows_data)))
    # Adjust body_h to exact sum based on each_h
    row_heights = [each_h] * n_rows_data

    # Make LEFT table adopt the compact row height so both sides align visually
    # Ensure enough rows already exist (done in _fill_tbl_gantt_left)
    for i in range(min(n_rows_data, len(left_tbl.rows) - 1)):
        left_tbl.rows[i + 1].height = each_h

    # Build compact grid in area
    grid = _build_time_grid_in_area(slide, "tblGanttTime", months, row_heights, left_tbl, gen_prefix)
    if not grid:
        return

    # Optional: draw today's line
    _draw_today_line(slide, grid, gen_prefix)

    # Draw bars (with smart labels)
    for idx, g in enumerate(gantt_rows):
        s = _to_date(_gget(g, "fecha_inicio"))
        e = _to_date(_gget(g, "fecha_fin"))
        if not s or not e:
            continue

        y = grid["row_tops"][idx]
        rh = grid["row_heights"][idx]
        bar_top = int(y + rh * 0.22)
        bar_h   = max(6, int(rh * 0.56))

        x0 = _date_to_x_grid(s, grid)
        x1 = _date_to_x_grid(e, grid)
        if x1 <= x0:
            x1 = x0 + 8

        width = int(x1 - x0)
        act_id = _norm_etapa_key(_gget(g, 'actividad', f'row_{idx+1}')) or f'row_{idx+1}'
        # Barra total (gris): duración completa
        total_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x0), int(bar_top), width, int(bar_h)
        )
        total_bar.name = f"{gen_prefix}bar_total_{act_id}"
        total_bar.fill.solid()
        total_bar.fill.fore_color.rgb = RGBColor(128, 128, 128)  # gris claro
        total_bar.line.color.rgb = RGBColor(0, 0, 0)

        # Barra de avance (encima): porcentaje de la columna "avance" (0..100)
        avance_pct = _to_percent_number(_gget(g, "avance", None))
        if avance_pct < 0:
            avance_pct = 0.0
        if avance_pct > 100:
            avance_pct = 100.0

        prog_w = int(width * (avance_pct / 100.0))
        if prog_w > 0:
            prog_w = max(2, prog_w)  # visible aunque sea muy pequeño
            prog = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(x0), int(bar_top), int(prog_w), int(bar_h)
            )
            prog.name = f"{gen_prefix}bar_prog_{act_id}"
            prog.fill.solid()

            col = (_gget(g, "color", None) or "VERDE").strip().upper()
            if col == "AMARILLO":
                prog.fill.fore_color.rgb = RGBColor(255, 192, 0)
            elif col == "ROJO":
                prog.fill.fore_color.rgb = RGBColor(192, 0, 0)
            else:
                prog.fill.fore_color.rgb = RGBColor(0, 176, 80)

            prog.line.color.rgb = RGBColor(0, 0, 0)

        # Globo de fecha (clonado desde shape oculto "globo_fecha")
        _add_globo_fecha(
            slide,
            x_end=int(x1),
            bar_top=int(bar_top),
            bar_h=int(bar_h),
            fecha_txt=_fmt_date(e),
            gen_prefix=gen_prefix,
            gap_px=12,
            scale_to_bar=0.88
        )







def fill_slide(slide, resumen, gantt_rows, riesgos_rows, etapas_row=None, gen_prefix="gen_gantt_"):
    shp = _find_shape(slide, "txt_titulo_folio_proyecto")
    if shp:
        _set_text_keep_style(shp, f"{resumen.folio_ppm} {resumen.nombre_proyecto}".strip())

    shp = _find_shape(slide, "txt_objetivo")
    if shp:
        _set_text_keep_style(shp, resumen.objetivo or "")

    shp = _find_shape(slide, "txt_descripcion_estatus")
    if shp:
        _set_descripcion_estatus_with_checks(
            slide,
            shp,
            raw_text=_gget(resumen, "descripcion_estatus", "") or "",
            marker="(*)",
        )


    tbl = _find_shape(slide, "tbl_kpi_horas_costos")
    if tbl and tbl.has_table:
        t = tbl.table
        donor = _get_table_donor_run(t)
        _set_cell_text_keep_style(t.cell(1, 0), resumen.horas_internas, donor)
        _set_cell_text_keep_style(t.cell(1, 1), resumen.horas_externas, donor)
        _set_cell_text_keep_style(t.cell(1, 2), resumen.horas_totales, donor)
        _set_cell_text_keep_style(t.cell(1, 3), _fmt_money(resumen.costo_total), donor)

    tbl = _find_shape(slide, "tbl_kpi_fechas")
    if tbl and tbl.has_table:
        t = tbl.table
        donor = _get_table_donor_run(t)
        _set_cell_text_keep_style(t.cell(1, 0), _fmt_date(resumen.fecha_inicio), donor)
        _set_cell_text_keep_style(t.cell(1, 1), _fmt_date(resumen.fecha_fin_liberacion), donor)
        _set_cell_text_keep_style(t.cell(1, 2), _fmt_date(resumen.fecha_fin_garantia), donor)

    tbl = _find_shape(slide, "tbl_kpi_avances")
    if tbl and tbl.has_table:
        t = tbl.table
        donor = _get_table_donor_run(t)
        _set_cell_text_keep_style(t.cell(1, 1), f"{round((_gget(resumen, 'avance_planeado', 0) or 0) * 100)}%", donor)
        _set_cell_text_keep_style(t.cell(1, 2), f"{round((_gget(resumen, 'avance_real', 0) or 0) * 100)}%", donor)

    sem = _find_shape(slide, "shp_estado_general")
    if sem:
        _color_estado_general_from_avances(
            sem,
            _gget(resumen, "avance_planeado", 0),
            _gget(resumen, "avance_real", 0),
        )

    tbl = _find_shape(slide, "tbl_responsables")
    if tbl and tbl.has_table:
        t = tbl.table
        donor = _get_table_donor_run(t)
        _set_cell_text_keep_style(t.cell(1, 0), resumen.direccion_area or "", donor)
        _set_cell_text_keep_style(t.cell(1, 1), resumen.lider_cliente or "", donor)
        _set_cell_text_keep_style(t.cell(1, 2), resumen.ern or "", donor)
        _set_cell_text_keep_style(t.cell(1, 3), resumen.le or "", donor)
        _set_cell_text_keep_style(t.cell(1, 4), resumen.ppm or "", donor)

    tbl = _find_shape(slide, "tbl_riesgos")
    if tbl and tbl.has_table:
        t = tbl.table
        donor = _get_table_donor_run(t)
        for r in range(1, len(t.rows)):
            for c in range(len(t.columns)):
                _set_cell_text_keep_style(t.cell(r, c), "", donor)
        max_rows = len(t.rows) - 1
        for i, rr in enumerate(riesgos_rows[:max_rows]):
            r = i + 1
            _set_cell_text_keep_style(t.cell(r, 0), _gget(rr, "folio_ppm", ""), donor)
            _set_cell_text_keep_style(t.cell(r, 1), _gget(rr, "riesgo", "") or "", donor)
            _set_cell_text_keep_style(t.cell(r, 2), _gget(rr, "responsable", "") or "", donor)
            _set_cell_text_keep_style(t.cell(r, 3), _gget(rr, "mitigacion", "") or "", donor)
            _set_cell_text_keep_style(t.cell(r, 4), _fmt_date(_gget(rr, "fecha_materializacion", None)), donor)

    # Left sidenote: ensure left has the rows created
    _fill_tbl_gantt_left(slide, gantt_rows, gen_prefix=gen_prefix)
    # Then draw compact right grid and bars
    _draw_gantt(slide, gantt_rows, gen_prefix=gen_prefix)

    # ETAPAS: iconos dentro de tbl_etapas_mantenimiento
    if etapas_row:
        set_etapas_icons_from_row(slide, etapas_row)




def build_etapas_by_folio(excel_path: str, sheet_name: str = "ETAPAS") -> dict:
    """
    Lee la hoja ETAPAS del Excel y devuelve un dict:
        { folio_ppm (str) : {col: value, ...} }
    Requiere que exista una columna 'folio_ppm'.
    """
    wb = load_workbook(excel_path, data_only=True)
    if sheet_name not in wb.sheetnames:
        raise ValueError(f"No existe la hoja '{sheet_name}' en el Excel.")
    ws = wb[sheet_name]

    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return {}

    headers = [str(h).strip() if h is not None else "" for h in rows[0]]
    try:
        folio_idx = [h.lower() for h in headers].index("folio_ppm")
    except ValueError:
        raise ValueError("La hoja ETAPAS debe tener una columna llamada 'folio_ppm'.")

    out = {}
    for r in rows[1:]:
        if r is None:
            continue
        folio = r[folio_idx]
        if folio is None or str(folio).strip() == "":
            continue
        folio_key = str(folio).strip()
        out[folio_key] = {headers[i]: r[i] for i in range(len(headers))}
    return out

def build_ppt_multi(template_path: str, out_path: str, projects, etapas_by_folio: dict = None, etapas_excel_path: str = None, etapas_sheet_name: str = "ETAPAS", gen_prefix="gen_gantt_"):
    """
    Genera un PPTX con 1 slide por proyecto a partir de la slide 0 de la plantilla.
    - projects puede ser una lista de tuplas:
        (resumen, gantt, riesgos)  o
        (resumen, gantt, riesgos, etapas_row)
    - Si se pasa etapas_by_folio (dict folio_ppm -> dict fila ETAPAS), y el item viene sin etapas_row,
      se resuelve automáticamente usando resumen.folio_ppm.
    """
    # AUTOLOAD_ETAPAS: si no te pasan etapas_by_folio, lo carga desde el Excel
    if etapas_by_folio is None and etapas_excel_path:
        etapas_by_folio = build_etapas_by_folio(etapas_excel_path, sheet_name=etapas_sheet_name)

    prs = Presentation(template_path)
    template_slide = prs.slides[0]

    for item in projects:
        if len(item) == 4:
            resumen, gantt, riesgos, etapas_row = item
        else:
            resumen, gantt, riesgos = item
            etapas_row = None
            if etapas_by_folio is not None:
                folio = str(getattr(resumen, "folio_ppm", "") or "").strip()
                etapas_row = etapas_by_folio.get(folio)

        s = duplicate_slide_full(prs, template_slide)
        fill_slide(s, resumen, gantt, riesgos, etapas_row=etapas_row, gen_prefix=gen_prefix)

    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    prs.save(out_path)





# ---- Helpers ETAPAS (normalización Excel) ----
def _strip_accents(s: str) -> str:
    import unicodedata
    s = "" if s is None else str(s)
    s = unicodedata.normalize("NFD", s)
    return "".join(ch for ch in s if unicodedata.category(ch) != "Mn")

def _norm_key_excel(s: str) -> str:
    import re
    s = _strip_accents(s).strip().upper()
    s = s.replace("\t", " ").replace("\n", " ").replace("\r", " ")
    s = re.sub(r"\s+", " ", s)
    return s.replace(" ", "_")

def _norm_status_excel(s) -> str:
    s = _strip_accents(s).strip().upper()
    s = s.replace("-", "_").replace(" ", "_")
    return s

def normalize_etapas_row(etapas_row: dict) -> dict:
    if not etapas_row:
        return {}
    out = {}
    for k, v in etapas_row.items():
        out[_norm_key_excel(k)] = v
    return out
# ---- End Helpers ETAPAS ----

# ===================== RESET ETAPAS (HOJA ETAPAS) =====================
# Fuente de iconos (en plantilla): ico_src_check, ico_src_gear (pueden estar ocultos)
# Tabla destino: tbl_etapas_mantenimiento
# Fila destino (blanca): ICON_ROW = 2

ETAPAS_COLS = [
    ("ESTIMACION", 0),
    ("PLANEACION", 1),
    ("ANALISIS_TECNICO", 2),
    ("DISENO_DETALLADO", 3),
    ("REALIZACION", 4),
    ("QA", 5),
    ("IMPLEMENTACION", 6),
    ("GARANTIA", 7),
]

def _strip_accents(s: str) -> str:
    s = str(s or "")
    s = unicodedata.normalize("NFD", s)
    return "".join(ch for ch in s if unicodedata.category(ch) != "Mn")

def _norm_key(s: str) -> str:
    s = _strip_accents(s).strip().upper()
    return s.replace(" ", "_")

def _norm_etapa_value(v) -> str:
    s = _strip_accents(v).strip().upper()
    return s.replace(" ", "_")

def normalize_etapas_row(etapas_row: dict) -> dict:
    """
    Convierte un dict de Excel (headers con acentos/mixtos) a llaves estandar:
      "Análisis_tecnico" -> "ANALISIS_TECNICO"
      "EN CURSO" -> "EN_CURSO"
    """
    out = {}
    for k, v in (etapas_row or {}).items():
        out[_norm_key(k)] = _norm_etapa_value(v)
    return out

def _find_shape_any_master_layout(slide, name: str):
    # slide
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    # layout
    try:
        for shp in slide.slide_layout.shapes:
            if shp.name == name:
                return shp
    except Exception:
        pass
    # master (python-pptx)
    try:
        for shp in slide.part.slide_master.shapes:
            if shp.name == name:
                return shp
    except Exception:
        pass
    return None

def _bring_to_front(slide, shape):
    try:
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.append(shape._element)
    except Exception:
        pass

def _remap_all_rids_to_slide_part(slide, icon_el):
    """
    Remapea TODOS los atributos r:embed / r:link encontrados en el XML del icono,
    incluyendo blip, svgBlip, extLst, etc.
    """
    src_parts = []
    try: src_parts.append(slide.part)
    except: pass
    try: src_parts.append(slide.slide_layout.part)
    except: pass
    try: src_parts.append(slide.part.slide_master.part)
    except: pass

    dst_part = slide.part
    rid_map = {}

    REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    EMBED = f"{{{REL_NS}}}embed"
    LINK  = f"{{{REL_NS}}}link"

    # recorre TODOS los nodos con atributos
    for el in icon_el.xpath(".//*[@*]"):
        for attr in (EMBED, LINK):
            old_rid = el.get(attr)
            if not old_rid:
                continue

            if old_rid in rid_map:
                el.set(attr, rid_map[old_rid])
                continue

            rel = None
            for sp in src_parts:
                try:
                    rel = sp.rels.get(old_rid)
                except Exception:
                    rel = None
                if rel is not None:
                    break

            if rel is None:
                continue

            new_rid = dst_part.relate_to(rel._target, rel.reltype)
            rid_map[old_rid] = new_rid
            el.set(attr, new_rid)


def set_etapas_icons_from_row(slide, etapas_row: dict):
    """
    Inserta (clona) los iconos ico_src_check / ico_src_gear dentro de la tabla
    tbl_etapas_mantenimiento, en la fila blanca (auto-detect), centrados en
    cada celda de etapa según el estatus leído del Excel.
    """
    tbl_shape = _find_shape_any_master_layout(slide, "tbl_etapas_mantenimiento")
    if not tbl_shape or not getattr(tbl_shape, "has_table", False):
        return

    # Normaliza llaves del Excel (quita acentos, espacios -> _)
    etapas_row = normalize_etapas_row(etapas_row)

    # limpiar iconos generados previamente
    for shp in list(slide.shapes):
        if getattr(shp, "name", "").startswith("gen_etapa_icon_"):
            el = shp._element
            el.getparent().remove(el)

    table = tbl_shape.table

    # Auto-detect fila destino:
    ICON_ROW = 2 if len(table.rows) >= 3 else 1

    # offsets columnas desde el SHAPE
    col_lefts = []
    acc_x = tbl_shape.left
    for c in table.columns:
        col_lefts.append(acc_x)
        acc_x += c.width

    # offsets filas desde el SHAPE
    row_tops = []
    acc_y = tbl_shape.top
    for r in table.rows:
        row_tops.append(acc_y)
        acc_y += r.height

    def _insert_icon_from_src(src_icon):
        icon_el = deepcopy(src_icon._element)
        _remap_all_rids_to_slide_part(slide, icon_el)
        slide.shapes._spTree.insert_element_before(icon_el, "p:extLst")
        return next((sh for sh in slide.shapes if sh._element is icon_el), None)

    for col_name, col_idx in ETAPAS_COLS:
        status = _norm_status_excel(etapas_row.get(col_name))

        # Solo insertar ícono cuando el estado es válido; en blanco u otro -> sin ícono
        if status == "COMPLETO":
            src_name = "ico_src_check"
        elif status == "EN_CURSO":
            src_name = "ico_src_gear"
        else:
            continue

        src_icon = _find_shape_any_master_layout(slide, src_name)
        if not src_icon:
            continue

        if col_idx >= len(table.columns) or ICON_ROW >= len(table.rows):
            continue

        cell_left = col_lefts[col_idx]
        cell_top = row_tops[ICON_ROW]
        cell_w = table.columns[col_idx].width
        cell_h = table.rows[ICON_ROW].height
        size = int(min(cell_w, cell_h) * 0.65)

        icon = _insert_icon_from_src(src_icon)
        if not icon:
            continue

        icon.width = size
        icon.height = size
        icon.left = int(cell_left + (cell_w - size) / 2)
        icon.top = int(cell_top + (cell_h - size) / 2)
        icon.name = f"gen_etapa_icon_{col_name}"
        try:
            icon.hidden = False
        except Exception:
            pass

        _force_visible(icon)
        _bring_to_front(slide, icon)

# =================== END RESET ETAPAS ===================